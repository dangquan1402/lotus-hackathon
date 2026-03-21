"""Generate Lumina Hackathon Pitch Deck (PPTX)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ────────────────────────────────────────────────────────────────────
DARK_GREEN = RGBColor(0x1E, 0x3A, 0x2F)
CREAM = RGBColor(0xFD, 0xF8, 0xF0)
GOLD = RGBColor(0xC8, 0x96, 0x3E)
MUTED = RGBColor(0x8A, 0x8A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0x2D, 0x4A, 0x3E)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
SOFT_GOLD = RGBColor(0xE8, 0xC8, 0x7A)

# ── Dimensions (16:9) ────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox, tf


def add_paragraph(tf, text, font_size=18, color=DARK_TEXT, bold=False,
                  font_name="Calibri", alignment=PP_ALIGN.LEFT,
                  space_before=0, space_after=0, line_spacing=1.2):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = Pt(int(font_size * line_spacing))
    return p


def add_accent_bar(slide, left, top, width=Inches(0.8), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bullet_color=GOLD, spacing=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
        p.line_spacing = Pt(int(font_size * 1.4))
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Arial")
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u25CF")  # filled circle
        # Set bullet color
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
        srgbClr.set("val", f"{bullet_color}")
    return txBox, tf


def make_section_title(slide, text, left=Inches(0.8), top=Inches(0.6)):
    """Add a section title with accent bar."""
    add_accent_bar(slide, left, top)
    add_text_box(slide, left, top + Inches(0.2), Inches(11), Inches(0.8),
                 text, font_size=36, color=DARK_GREEN, bold=True,
                 font_name="Calibri")


def create_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN)

    # Diamond icon
    add_text_box(slide, Inches(5.7), Inches(1.5), Inches(2), Inches(1),
                 "\U0001F48E", font_size=64, color=GOLD, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.2),
                 "Lumina", font_size=72, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # Subtitle line 1
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.8),
                 "AI-Powered Personalized Learning", font_size=28,
                 color=CREAM, alignment=PP_ALIGN.CENTER)

    # Accent bar
    add_accent_bar(slide, Inches(5.5), Inches(4.7), Inches(2.3), Inches(0.04))

    # Tagline
    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.6),
                 "Learn anything, your way", font_size=22,
                 color=SOFT_GOLD, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — The Problem
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, CREAM)
    make_section_title(slide, "The Problem")

    problems = [
        "One-size-fits-all education doesn't work",
        "Learners have different styles, levels, and goals",
        "Content is static, not adaptive to individual needs",
        "No way to assess what you already know before diving in",
    ]
    add_bullet_list(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(4.5),
                    problems, font_size=22, color=DARK_TEXT, spacing=12)

    # Bottom accent
    add_text_box(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.6),
                 "Every learner deserves content that meets them where they are.",
                 font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — Our Solution
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, CREAM)
    make_section_title(slide, "Our Solution")

    add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.6),
                 "Lumina creates personalized, multi-modal learning experiences:",
                 font_size=20, color=MUTED)

    solutions = [
        "Pre-assessment chat to gauge existing knowledge before generating content",
        "AI-generated lessons tailored to your level, style, and interests",
        "Videos, slides, quizzes, and flashcards — all auto-generated from one topic",
        "AI Chat Tutor — streaming conversational tutor grounded in lesson content",
        "Animated video with Grok — static images converted to 6s animated clips",
        "Knowledge graph tracking concepts learned across sessions",
        "Study plan with AI-recommended next topics based on your progress",
    ]
    add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(11), Inches(4.5),
                    solutions, font_size=20, color=DARK_TEXT, spacing=10)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — How It Works (Technical)
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, CREAM)
    make_section_title(slide, "How It Works")

    # Flow diagram as text boxes with arrows
    flow_steps = [
        ("User Profile", DARK_GREEN),
        ("\u2192", GOLD),
        ("Pre-Assessment\n(LangGraph)", DARK_GREEN),
        ("\u2192", GOLD),
        ("Web Search\n(Firecrawl)", DARK_GREEN),
        ("\u2192", GOLD),
        ("Content Synthesis\n(LLM)", DARK_GREEN),
        ("\u2192", GOLD),
        ("Grok Imagine\n(Animated Video)", DARK_GREEN),
        ("\u2192", GOLD),
        ("Multi-modal Output\n(Video/Slides/Quiz)", DARK_GREEN),
    ]

    x_start = Inches(0.2)
    box_width = Inches(1.8)
    arrow_width = Inches(0.35)
    y_pos = Inches(2.0)

    for i, (text, color) in enumerate(flow_steps):
        if text == "\u2192":
            add_text_box(slide, x_start, y_pos + Inches(0.15), arrow_width, Inches(0.6),
                         "\u2192", font_size=28, color=GOLD, bold=True,
                         alignment=PP_ALIGN.CENTER)
            x_start += arrow_width
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x_start, y_pos, box_width, Inches(0.9))
            shape.fill.solid()
            shape.fill.fore_color.rgb = DARK_GREEN
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(11)
            p.font.color.rgb = CREAM
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            x_start += box_width

    # Tech stack section
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
                 "Tech Stack", font_size=20, color=GOLD, bold=True)

    stack_items = [
        "Frontend: React + Vite + Tailwind CSS",
        "Backend: FastAPI + SQLAlchemy + PostgreSQL",
        "AI Pipeline: LangGraph (assessment + chat) \u2192 FuseAPI (multi-model LLM routing) \u2192 Remotion (video)",
        "Search: Firecrawl (web scraping + search)",
        "Voice: mlx-whisper (alignment), TTS narration",
        "Visualization: Markmap (mind maps), react-force-graph (knowledge graph)",
    ]
    add_bullet_list(slide, Inches(1.0), Inches(4.1), Inches(11), Inches(3),
                    stack_items, font_size=16, color=DARK_TEXT, spacing=6)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Key Features
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, CREAM)
    make_section_title(slide, "Key Features")

    features = [
        ("Chat-based Pre-Assessment", "LangGraph-powered conversation to gauge knowledge before lesson generation"),
        ("AI Chat Tutor with Streaming", "Real-time SSE conversational tutor grounded in lesson content + web search"),
        ("Animated Video Generation", "Grok Imagine converts images to 6s clips, Remotion renders with timing sync"),
        ("Playbook Mode", "Per-section audio with auto-play and auto-advance in slide presenter"),
        ("Interactive Slide Presenter", "In-browser slides with AI images, keyboard nav, and PPTX download"),
        ("7 Visual Styles", "Cartoon, watercolor, photorealistic, minimalist, anime, scientific, 3D render"),
        ("Cross-session Knowledge Graph", "Force-graph linking topics to concepts across all sessions"),
        ("Mind Map & Flashcards", "Markmap concept trees + flip cards with mastery tracking"),
        ("Learning Analytics & Study Plan", "Timeline, session stats, and AI-recommended next topics"),
    ]

    # Two-column layout
    col_left = Inches(0.8)
    col_right = Inches(6.8)
    y = Inches(1.6)
    row_height = Inches(0.6)

    for i, (title, desc) in enumerate(features):
        col = col_left if i < 5 else col_right
        row_y = y + (i % 5) * row_height

        # Number badge
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, col, row_y + Inches(0.05), Inches(0.35), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GOLD
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Feature title + description
        _, tf = add_text_box(slide, col + Inches(0.5), row_y, Inches(5.2), Inches(0.55),
                             title, font_size=15, color=DARK_GREEN, bold=True)
        add_paragraph(tf, desc, font_size=11, color=MUTED, line_spacing=1.1)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Demo Flow
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "DEMO FLOW", font_size=14, color=GOLD, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.0))
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                 "Live Walkthrough", font_size=36, color=CREAM, bold=True)

    demo_steps = [
        ("1", "Enter topic", "Type any learning topic into the search bar"),
        ("2", "Pre-assessment chat", "AI asks targeted questions to gauge existing knowledge"),
        ("3", "AI generates lesson", "Web search + LLM synthesizes personalized content"),
        ("4", "View slides with images", "Navigate AI-generated slides with multiple images per section"),
        ("5", "Chat with AI tutor", "Ask follow-up questions with streaming responses and source citations"),
        ("6", "Watch animated video", "AI-generated animated clips with narrated audio"),
        ("7", "Take quiz", "Test understanding with auto-generated questions"),
        ("8", "Explore mind map", "Interactive concept tree visualization"),
        ("9", "Review knowledge graph", "See connections across all learning sessions"),
    ]

    y_start = Inches(2.0)
    for i, (num, title, desc) in enumerate(demo_steps):
        y_pos = y_start + i * Inches(0.57)

        # Step number
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.0), y_pos, Inches(0.35), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GOLD
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = DARK_GREEN
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Title
        add_text_box(slide, Inches(1.6), y_pos + Inches(0.02), Inches(3.5), Inches(0.4),
                     title, font_size=18, color=CREAM, bold=True)
        # Description
        add_text_box(slide, Inches(5.2), y_pos + Inches(0.02), Inches(7), Inches(0.4),
                     desc, font_size=15, color=SOFT_GOLD)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — What Makes Us Different
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, CREAM)
    make_section_title(slide, "What Makes Us Different")

    differentiators = [
        ("Pre-assessment BEFORE content generation",
         "Most platforms test you after — Lumina adapts the lesson itself based on what you already know"),
        ("Multi-modal from ONE topic",
         "Video + slides + quiz + flashcards + mind map — all generated from a single topic input"),
        ("Real-time streaming AI chat",
         "LangGraph-powered conversational tutor with SSE streaming and source citations from web search"),
        ("Cross-session knowledge graph",
         "Obsidian-style concept linking across all your learning sessions"),
        ("Open architecture",
         "FuseAPI routes to any LLM (Gemini, Claude, GPT). Grok for animation. Swap any component."),
    ]

    y = Inches(1.8)
    for i, (title, desc) in enumerate(differentiators):
        # Gold accent line
        add_accent_bar(slide, Inches(1.0), y + Inches(0.15), Inches(0.05), Inches(0.5))

        add_text_box(slide, Inches(1.3), y, Inches(10.5), Inches(0.4),
                     title, font_size=18, color=DARK_GREEN, bold=True)
        add_text_box(slide, Inches(1.3), y + Inches(0.4), Inches(10.5), Inches(0.45),
                     desc, font_size=14, color=MUTED)
        y += Inches(1.0)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Tech Stack
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "ARCHITECTURE", font_size=14, color=GOLD, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.0))
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
                 "Tech Stack", font_size=36, color=CREAM, bold=True)

    stack = [
        ("Frontend", "React, Vite, Tailwind CSS"),
        ("Backend", "FastAPI, SQLAlchemy, PostgreSQL"),
        ("AI / LLM", "LangGraph (assessment + chat), FuseAPI (multi-model routing)"),
        ("Search", "Firecrawl (web scraping + search)"),
        ("Video", "Remotion (programmatic video generation)"),
        ("Animation", "Grok Imagine API (image-to-video)"),
        ("Voice", "mlx-whisper (alignment), TTS narration"),
        ("Visualization", "Markmap (mind maps), react-force-graph (knowledge graph)"),
        ("Infrastructure", "Docker Compose (4 services: frontend, backend, db, remotion)"),
    ]

    # Two-column card layout
    col1_x = Inches(0.8)
    col2_x = Inches(6.8)
    y = Inches(2.2)

    for i, (label, tech) in enumerate(stack):
        col = col1_x if i < 5 else col2_x
        row_y = y + (i % 5) * Inches(0.95)

        # Card background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, col, row_y, Inches(5.5), Inches(0.78))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GREEN
        shape.line.fill.background()

        # Label
        add_text_box(slide, col + Inches(0.3), row_y + Inches(0.1), Inches(4.8), Inches(0.3),
                     label, font_size=12, color=GOLD, bold=True)
        # Tech
        add_text_box(slide, col + Inches(0.3), row_y + Inches(0.45), Inches(4.8), Inches(0.4),
                     tech, font_size=15, color=CREAM)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Future Vision
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, CREAM)
    make_section_title(slide, "Future Vision")

    vision_items = [
        "Real-time collaborative learning (study groups)",
        "Spaced repetition algorithm for flashcard scheduling",
        "Teacher / instructor mode with class analytics",
        "Mobile app (Capacitor integration already scaffolded)",
        "Multi-language support with localized narration",
        "Adaptive difficulty that adjusts in real-time during lessons",
    ]
    add_bullet_list(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(4.5),
                    vision_items, font_size=22, color=DARK_TEXT, spacing=14)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Thank You
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN)

    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1),
                 "\U0001F48E", font_size=64, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1),
                 "Thank You", font_size=56, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, Inches(5.5), Inches(3.7), Inches(2.3), Inches(0.04))

    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.6),
                 "Built at FuseAPI AI Hackathon 2025", font_size=22,
                 color=CREAM, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
                 "Try it: localhost:3000", font_size=18,
                 color=SOFT_GOLD, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
                 "Lumina \u2014 Learn anything, your way",
                 font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    output_path = "Lumina_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Pitch deck saved to {output_path}")


if __name__ == "__main__":
    create_deck()
