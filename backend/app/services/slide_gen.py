from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Color palette ──────────────────────────────────────────────────────────────
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_LIGHT = RGBColor(0x81, 0x7A, 0xF3)
DARK_BG = RGBColor(0x0F, 0x0F, 0x1A)
CARD_BG = RGBColor(0x1A, 0x1A, 0x2E)
CARD_BG_LIGHT = RGBColor(0x22, 0x22, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
GREEN_DARK = RGBColor(0x16, 0x6A, 0x34)
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
SLIDE_W = Inches(16)
SLIDE_H = Inches(9)

# ── Reusable layout constants ─────────────────────────────────────────────────
MARGIN_X = Inches(0.8)
MARGIN_Y = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN_X


# ── Helper functions ──────────────────────────────────────────────────────────


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """Add a rectangle shape with optional border."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rounded rectangle (shape type 5)."""
    shape = slide.shapes.add_shape(5, left, top, width, height)  # 5 = rounded rect
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    bold=False,
    color=WHITE,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
    line_spacing=1.2,
):
    """Add a text box with full formatting control."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    pf = p._pPr
    if pf is not None:
        from pptx.oxml.ns import qn

        lnSpc = pf.makeelement(qn("a:lnSpc"), {})
        spcPct = lnSpc.makeelement(qn("a:spcPct"), {"val": str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pf.append(lnSpc)

    return tf


def _add_accent_bar(slide, left, top, width, height=Inches(0.06)):
    """Add a gradient-like accent bar."""
    _add_shape(slide, left, top, width, height, INDIGO)


def _add_image_safe(slide, img_path: Path, left, top, width, height):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), left, top, width, height)
        return True
    return False


def _add_page_number(slide, number: int):
    """Add a subtle page number at bottom-right."""
    _add_text(
        slide,
        SLIDE_W - Inches(1.2),
        SLIDE_H - Inches(0.5),
        Inches(0.8),
        Inches(0.4),
        str(number),
        font_size=11,
        color=MUTED,
        alignment=PP_ALIGN.RIGHT,
    )


# ── Slide builders ────────────────────────────────────────────────────────────


def _make_title_slide(prs, content: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BG)

    # Top accent bar (full width)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07))

    # Decorative side accent
    _add_shape(slide, Inches(0), Inches(0), Inches(0.07), SLIDE_H, INDIGO)

    # "LESSON" label
    _add_text(
        slide,
        Inches(2),
        Inches(2.2),
        Inches(12),
        Inches(0.5),
        "LESSON",
        font_size=14,
        bold=True,
        color=INDIGO_LIGHT,
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri",
    )

    # Title
    _add_text(
        slide,
        Inches(1.5),
        Inches(2.8),
        Inches(13),
        Inches(2.2),
        content["title"],
        font_size=48,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri Light",
        line_spacing=1.1,
    )

    # Divider line under title
    _add_shape(slide, Inches(6.5), Inches(5.1), Inches(3), Inches(0.04), INDIGO)

    # Overview subtitle
    overview = content.get("overview", "")
    if overview:
        _add_text(
            slide,
            Inches(2.5),
            Inches(5.5),
            Inches(11),
            Inches(1.5),
            overview,
            font_size=20,
            color=MUTED,
            alignment=PP_ALIGN.CENTER,
            line_spacing=1.4,
        )

    # Section count badge at bottom
    num_sections = len(content.get("sections", []))
    num_quiz = len(content.get("quiz_questions", []))
    badge_parts = [f"{num_sections} sections"]
    if num_quiz:
        badge_parts.append(f"{num_quiz} quiz questions")
    badge_text = "  |  ".join(badge_parts)

    # Badge background
    badge_w = Inches(6)
    badge_x = (SLIDE_W - badge_w) // 2
    _add_rounded_rect(slide, badge_x, Inches(7.4), badge_w, Inches(0.6), CARD_BG, INDIGO)
    _add_text(
        slide,
        badge_x,
        Inches(7.48),
        badge_w,
        Inches(0.5),
        badge_text,
        font_size=14,
        color=LIGHT_GRAY,
        alignment=PP_ALIGN.CENTER,
    )


def _make_section_slide(prs, i: int, section: dict, total: int, images_dir: Path, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BG)

    # Find all images for this section
    image_files = sorted(images_dir.glob(f"scene_{i:02d}_*.jpg"))
    if not image_files:
        image_files = [images_dir / f"scene_{i:02d}.jpg"]
    existing_images = [f for f in image_files if f.exists()]
    num_images = len(existing_images)

    # Top accent bar
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W)

    # Section number badge (styled pill)
    badge_w = Inches(1.6)
    _add_rounded_rect(slide, MARGIN_X, Inches(0.5), badge_w, Inches(0.5), INDIGO)
    _add_text(
        slide,
        MARGIN_X,
        Inches(0.53),
        badge_w,
        Inches(0.45),
        f"Section {i + 1} of {total}",
        font_size=13,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Section title
    title_max_w = Inches(8.5) if num_images > 0 else Inches(14)
    _add_text(
        slide,
        MARGIN_X,
        Inches(1.3),
        title_max_w,
        Inches(1),
        section["title"],
        font_size=34,
        bold=True,
        color=WHITE,
        font_name="Calibri Light",
        line_spacing=1.1,
    )

    # Accent line under title
    _add_shape(slide, MARGIN_X, Inches(2.35), Inches(2.5), Inches(0.05), INDIGO)

    # ── Layout depends on image count ──
    if num_images == 0:
        # Full-width text
        _add_text(
            slide,
            MARGIN_X,
            Inches(2.8),
            Inches(14),
            Inches(5.5),
            section["narration_text"],
            font_size=19,
            color=OFF_WHITE,
            line_spacing=1.5,
        )
    elif num_images == 1:
        # Left text, right single large image
        _add_text(
            slide,
            MARGIN_X,
            Inches(2.8),
            Inches(8.2),
            Inches(5.5),
            section["narration_text"],
            font_size=18,
            color=OFF_WHITE,
            line_spacing=1.5,
        )
        img_x = Inches(9.5)
        img_w = Inches(6)
        img_h = Inches(6)
        _add_image_safe(slide, existing_images[0], img_x, Inches(1.3), img_w, img_h)
    elif num_images == 2:
        # Text on top-left, two images stacked on right
        _add_text(
            slide,
            MARGIN_X,
            Inches(2.8),
            Inches(8.2),
            Inches(5.5),
            section["narration_text"],
            font_size=18,
            color=OFF_WHITE,
            line_spacing=1.5,
        )
        img_x = Inches(9.5)
        img_w = Inches(6)
        img_h = Inches(3.5)
        gap = Inches(0.3)
        _add_image_safe(slide, existing_images[0], img_x, Inches(0.5), img_w, img_h)
        _add_image_safe(slide, existing_images[1], img_x, Inches(0.5) + img_h + gap, img_w, img_h)
    else:
        # 3+ images: text left, images in a grid on right (2 cols)
        _add_text(
            slide,
            MARGIN_X,
            Inches(2.8),
            Inches(7.5),
            Inches(5.5),
            section["narration_text"],
            font_size=17,
            color=OFF_WHITE,
            line_spacing=1.5,
        )
        img_x = Inches(9.0)
        col_w = Inches(3.3)
        row_h = Inches(2.5)
        gap = Inches(0.2)
        for idx, img_file in enumerate(existing_images[:6]):
            col = idx % 2
            row = idx // 2
            x = img_x + col * (col_w + gap)
            y = Inches(0.5) + row * (row_h + gap)
            _add_image_safe(slide, img_file, x, y, col_w, row_h)

    _add_page_number(slide, page_num)


def _make_quiz_title_slide(prs, num_questions: int, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)

    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W)

    # Quiz icon area
    _add_rounded_rect(
        slide,
        Inches(6.75),
        Inches(2.2),
        Inches(2.5),
        Inches(0.7),
        INDIGO,
    )
    _add_text(
        slide,
        Inches(6.75),
        Inches(2.27),
        Inches(2.5),
        Inches(0.6),
        "QUIZ TIME",
        font_size=18,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    _add_text(
        slide,
        Inches(2),
        Inches(3.5),
        Inches(12),
        Inches(1.5),
        "Knowledge Check",
        font_size=48,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri Light",
    )

    _add_shape(slide, Inches(6.5), Inches(5.2), Inches(3), Inches(0.04), INDIGO)

    _add_text(
        slide,
        Inches(3),
        Inches(5.6),
        Inches(10),
        Inches(1),
        f"Test your understanding with {num_questions} questions",
        font_size=20,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    _add_page_number(slide, page_num)


def _make_quiz_question_slide(prs, qi: int, q: dict, total_q: int, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)

    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W)

    # Question number badge
    badge_w = Inches(2.2)
    _add_rounded_rect(slide, MARGIN_X, Inches(0.5), badge_w, Inches(0.5), INDIGO)
    _add_text(
        slide,
        MARGIN_X,
        Inches(0.53),
        badge_w,
        Inches(0.45),
        f"Question {qi + 1} of {total_q}",
        font_size=13,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
    )

    # Question text
    _add_text(
        slide,
        MARGIN_X,
        Inches(1.5),
        Inches(14),
        Inches(1.5),
        q["question"],
        font_size=26,
        bold=True,
        color=WHITE,
        line_spacing=1.3,
    )

    # Options — two-column layout for 4 options, single column otherwise
    options = q.get("options", [])
    correct_idx = q.get("correct_index", -1)

    if len(options) == 4:
        # 2x2 grid layout
        card_w = Inches(6.8)
        card_h = Inches(1.2)
        gap_x = Inches(0.4)
        gap_y = Inches(0.3)
        start_x = Inches(0.8)
        start_y = Inches(3.5)

        for oi, opt in enumerate(options):
            col = oi % 2
            row = oi // 2
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)
            letter = chr(65 + oi)
            is_correct = oi == correct_idx

            bg_color = GREEN_DARK if is_correct else CARD_BG
            border_color = GREEN if is_correct else RGBColor(0x33, 0x33, 0x55)
            _add_rounded_rect(slide, x, y, card_w, card_h, bg_color, border_color)

            # Letter badge
            letter_w = Inches(0.5)
            _add_rounded_rect(
                slide,
                x + Inches(0.2),
                y + Inches(0.35),
                letter_w,
                Inches(0.5),
                INDIGO if not is_correct else GREEN,
            )
            _add_text(
                slide,
                x + Inches(0.2),
                y + Inches(0.38),
                letter_w,
                Inches(0.45),
                letter,
                font_size=16,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
            )

            # Option text
            _add_text(
                slide,
                x + Inches(0.9),
                y + Inches(0.25),
                card_w - Inches(1.1),
                card_h - Inches(0.3),
                opt,
                font_size=16,
                color=WHITE if is_correct else OFF_WHITE,
                bold=is_correct,
                line_spacing=1.2,
            )
    else:
        # Single column layout
        option_y = Inches(3.5)
        card_w = Inches(13)
        card_h = Inches(0.9)
        card_x = Inches(1.5)

        for oi, opt in enumerate(options):
            letter = chr(65 + oi)
            is_correct = oi == correct_idx

            bg_color = GREEN_DARK if is_correct else CARD_BG
            border_color = GREEN if is_correct else RGBColor(0x33, 0x33, 0x55)
            _add_rounded_rect(slide, card_x, option_y, card_w, card_h, bg_color, border_color)

            # Letter badge
            letter_w = Inches(0.5)
            _add_rounded_rect(
                slide,
                card_x + Inches(0.2),
                option_y + Inches(0.2),
                letter_w,
                Inches(0.5),
                INDIGO if not is_correct else GREEN,
            )
            _add_text(
                slide,
                card_x + Inches(0.2),
                option_y + Inches(0.23),
                letter_w,
                Inches(0.45),
                letter,
                font_size=16,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
            )

            _add_text(
                slide,
                card_x + Inches(0.9),
                option_y + Inches(0.2),
                card_w - Inches(1.2),
                Inches(0.6),
                opt,
                font_size=17,
                color=WHITE if is_correct else OFF_WHITE,
                bold=is_correct,
            )

            option_y += card_h + Inches(0.2)

    # Explanation at bottom in a card
    explanation = q.get("explanation", "")
    if explanation:
        exp_y = Inches(6.8)
        exp_w = Inches(13.5)
        exp_x = Inches(1.2)
        _add_rounded_rect(slide, exp_x, exp_y, exp_w, Inches(1.4), CARD_BG_LIGHT)
        _add_shape(slide, exp_x, exp_y, Inches(0.06), Inches(1.4), AMBER)
        _add_text(
            slide,
            exp_x + Inches(0.4),
            exp_y + Inches(0.15),
            exp_w - Inches(0.6),
            Inches(1.1),
            explanation,
            font_size=14,
            color=MUTED,
            line_spacing=1.3,
        )

    _add_page_number(slide, page_num)


def _make_summary_slide(prs, content: dict, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)

    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W)

    # "RECAP" label
    _add_text(
        slide,
        Inches(2),
        Inches(0.8),
        Inches(12),
        Inches(0.5),
        "RECAP",
        font_size=14,
        bold=True,
        color=INDIGO_LIGHT,
        alignment=PP_ALIGN.CENTER,
    )

    _add_text(
        slide,
        Inches(1),
        Inches(1.3),
        Inches(14),
        Inches(1.2),
        "Key Takeaways",
        font_size=42,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
        font_name="Calibri Light",
    )

    _add_shape(slide, Inches(6.5), Inches(2.7), Inches(3), Inches(0.04), INDIGO)

    # Section cards in a grid
    sections = content.get("sections", [])
    num_sections = len(sections)

    if num_sections <= 4:
        # Single column centered cards
        card_w = Inches(10)
        card_h = Inches(0.9)
        card_x = Inches(3)
        start_y = Inches(3.3)
        gap = Inches(0.25)

        for idx, section in enumerate(sections):
            y = start_y + idx * (card_h + gap)
            _add_rounded_rect(slide, card_x, y, card_w, card_h, CARD_BG)

            # Number circle
            circle_w = Inches(0.6)
            _add_rounded_rect(
                slide, card_x + Inches(0.2), y + Inches(0.15), circle_w, circle_w, INDIGO
            )
            _add_text(
                slide,
                card_x + Inches(0.2),
                y + Inches(0.2),
                circle_w,
                Inches(0.5),
                str(idx + 1),
                font_size=18,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
            )

            # Section title
            _add_text(
                slide,
                card_x + Inches(1.1),
                y + Inches(0.2),
                card_w - Inches(1.5),
                Inches(0.6),
                section["title"],
                font_size=20,
                color=OFF_WHITE,
            )
    else:
        # Two-column layout for 5+ sections
        card_w = Inches(6.5)
        card_h = Inches(0.8)
        gap_x = Inches(0.5)
        gap_y = Inches(0.2)
        start_x = Inches(1.2)
        start_y = Inches(3.3)

        for idx, section in enumerate(sections):
            col = idx % 2
            row = idx // 2
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)

            _add_rounded_rect(slide, x, y, card_w, card_h, CARD_BG)

            # Number
            circle_w = Inches(0.5)
            _add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), circle_w, circle_w, INDIGO)
            _add_text(
                slide,
                x + Inches(0.15),
                y + Inches(0.18),
                circle_w,
                Inches(0.45),
                str(idx + 1),
                font_size=16,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
            )

            _add_text(
                slide,
                x + Inches(0.85),
                y + Inches(0.18),
                card_w - Inches(1.1),
                Inches(0.5),
                section["title"],
                font_size=17,
                color=OFF_WHITE,
            )

    # Bottom message
    _add_text(
        slide,
        Inches(3),
        Inches(8),
        Inches(10),
        Inches(0.6),
        "Great job completing this lesson!",
        font_size=18,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    _add_page_number(slide, page_num)


# ── Main entry point ──────────────────────────────────────────────────────────


async def agenerate_slides(
    session_id: int,
    content: dict,
    images_dir: Path,
    output_dir: Path,
) -> Path:
    """Generate a polished PPTX presentation from learning content."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    sections = content.get("sections", [])
    quiz_questions = content.get("quiz_questions", [])
    page = 1

    # 1. Title slide
    _make_title_slide(prs, content)
    page += 1

    # 2. Content slides (one per section)
    for i, section in enumerate(sections):
        _make_section_slide(prs, i, section, len(sections), images_dir, page)
        page += 1

    # 3. Summary/recap slide
    _make_summary_slide(prs, content, page)
    page += 1

    # 4. Quiz slides
    if quiz_questions:
        _make_quiz_title_slide(prs, len(quiz_questions), page)
        page += 1
        for qi, q in enumerate(quiz_questions):
            _make_quiz_question_slide(prs, qi, q, len(quiz_questions), page)
            page += 1

    output_path = output_dir / f"lesson_{session_id}.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
